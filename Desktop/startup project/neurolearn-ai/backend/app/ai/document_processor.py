import fitz
from docx import Document
from pptx import Presentation
import pytesseract
from PIL import Image
import io
import os
from typing import Optional, List
from app.config import settings
from loguru import logger


class DocumentProcessor:
    def __init__(self):
        if settings.ocr_enabled:
            pytesseract.pytesseract.tesseract_cmd = settings.tesseract_cmd

    async def extract_text(self, file_path: str, file_type: str) -> str:
        try:
            if file_type == "pdf":
                return await self._extract_from_pdf(file_path)
            elif file_type == "docx":
                return await self._extract_from_docx(file_path)
            elif file_type == "pptx":
                return await self._extract_from_pptx(file_path)
            elif file_type in ["txt", "md"]:
                return await self._extract_from_text(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
        except Exception as e:
            logger.error(f"Text extraction failed for {file_path}: {e}")
            raise

    async def _extract_from_pdf(self, file_path: str) -> str:
        text = ""
        doc = fitz.open(file_path)
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            page_text = page.get_text()
            text += page_text + "\n"
            
            if settings.ocr_enabled:
                pix = page.get_pixmap()
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                ocr_text = pytesseract.image_to_string(img)
                if ocr_text.strip():
                    text += ocr_text + "\n"
        
        doc.close()
        return text.strip()

    async def _extract_from_docx(self, file_path: str) -> str:
        doc = Document(file_path)
        text = ""
        
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " "
                text += "\n"
        
        return text.strip()

    async def _extract_from_pptx(self, file_path: str) -> str:
        prs = Presentation(file_path)
        text = ""
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        
        return text.strip()

    async def _extract_from_text(self, file_path: str) -> str:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()

    def clean_text(self, text: str) -> str:
        text = text.strip()
        text = "\n".join(line.strip() for line in text.split("\n") if line.strip())
        text = " ".join(text.split())
        return text

    def chunk_text(self, text: str, chunk_size: int = None, overlap: int = None) -> List[str]:
        chunk_size = chunk_size or settings.chunk_size
        overlap = overlap or settings.chunk_overlap
        
        chunks = []
        start = 0
        text_length = len(text)
        
        while start < text_length:
            end = start + chunk_size
            chunk = text[start:end]
            
            if end < text_length:
                last_space = chunk.rfind(" ")
                if last_space > chunk_size // 2:
                    chunk = chunk[:last_space]
                    end = start + last_space
            
            chunks.append(chunk.strip())
            start = end - overlap
        
        return [chunk for chunk in chunks if chunk]
